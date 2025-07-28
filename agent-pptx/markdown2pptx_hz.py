import re
import chardet
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def safe_read_file(file_path):
    """安全读取文件，自动检测编码"""
    try:
        # 方法1: 使用chardet检测编码
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            detected_encoding = result['encoding']
            confidence = result['confidence']
            
        print(f"检测到的编码: {detected_encoding} (置信度: {confidence:.2f})")
        
        # 如果置信度较高，使用检测到的编码
        if confidence > 0.7:
            with open(file_path, 'r', encoding=detected_encoding) as f:
                return f.read()
    except Exception as e:
        print(f"编码检测失败: {e}")
    
    # 方法2: 尝试常见编码
    encodings = ['utf-8', 'gbk', 'gb2312', 'cp936', 'latin-1', 'cp1252']
    
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                content = f.read()
                print(f"成功使用编码: {encoding}")
                return content
        except UnicodeDecodeError:
            continue
        except Exception as e:
            print(f"使用编码 {encoding} 时出错: {e}")
            continue
    
    # 方法3: 最后的fallback - 使用utf-8并忽略错误
    try:
        print("使用UTF-8编码并忽略错误字符")
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        print(f"最终读取失败: {e}")
        raise ValueError(f"无法读取文件 {file_path}，请检查文件是否存在或已损坏")


def parse_markdown(file_path):
    """解析Markdown文件"""
    # 使用安全读取方法
    content = safe_read_file(file_path)
    
    slides_data = []
    current_slide = None
    
    lines = content.split('\n')
    i = 0
    
    while i < len(lines):
        line = lines[i].strip()
        
        # 一级标题 - 新幻灯片
        if line.startswith('# '):
            if current_slide:
                slides_data.append(current_slide)
            
            current_slide = {
                'title': line[2:].strip(),
                'description': '',
                'modules': []
            }
            i += 1
            
            # 收集描述文字（一级标题后到二级标题前的内容）
            description_lines = []
            while i < len(lines):
                if lines[i].strip().startswith('## '):
                    break
                if lines[i].strip():
                    description_lines.append(lines[i].strip())
                i += 1
            
            current_slide['description'] = '\n'.join(description_lines)  # 描述也保持换行
            continue
            
        # 二级标题 - 模块标题
        elif line.startswith('## '):
            if current_slide:
                module_title = line[3:].strip()
                module_content = []
                i += 1
                
                # 收集该模块的内容（直到下一个二级标题或一级标题）
                while i < len(lines):
                    next_line = lines[i].strip()
                    if next_line.startswith('# ') or next_line.startswith('## '):
                        break
                    if next_line:
                        module_content.append(next_line)
                    elif lines[i] == '':  # 保留空行
                        module_content.append('')
                    i += 1
                
                current_slide['modules'].append({
                    'title': module_title,
                    'content': '\n'.join(module_content)  # 使用换行符连接，保持结构
                })
                continue
        
        i += 1
    
    # 添加最后一个幻灯片
    if current_slide:
        slides_data.append(current_slide)
    
    return slides_data


def apply_text_formatting(text_frame, text_content):
    """应用文本格式，处理粗体标记，保持换行结构"""
    # 清空现有内容
    text_frame.clear()
    
    # 按行分割文本内容
    lines = text_content.split('\n')
    
    for line_idx, line in enumerate(lines):
        line = line.strip()
        if not line:  # 跳过空行，但保持段落分隔
            if line_idx < len(lines) - 1:  # 不是最后一行
                text_frame.add_paragraph()
            continue
            
        # 为每一行创建段落（第一行使用默认段落）
        if line_idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # 使用正则表达式分割该行文本，保留**标记
        parts = re.split(r'(\*\*.*?\*\*)', line)
        
        # 在同一段落中处理加粗和普通文字
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                # 粗体文本 - 在同一行内，不换行
                run = p.add_run()
                run.text = part[2:-2]  # 去掉**标记
                run.font.bold = True
            else:
                # 普通文本
                if part:  # 非空字符串
                    run = p.add_run()
                    run.text = part


def create_pptx_from_markdown(markdown_file, output_file):
    """从Markdown创建PPTX"""
    # 解析Markdown
    slides_data = parse_markdown(markdown_file)
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(13.33)  # 16:9比例
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides_data:
        # 创建空白幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 1. 顶层主标题
        title_left = Cm(0)
        title_top = Cm(0)
        title_width = Cm(30)
        title_height = Cm(1.5)
        
        title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_shape.text_frame
        title_frame.text = slide_data['title']
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title_frame.paragraphs[0].font.name = '微软雅黑'
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.color.rgb = RGBColor(139, 0, 0)  # 深红色
        
        # 2. 中层矩形解释框
        desc_left = Cm(0)
        desc_top = Cm(1.7)
        desc_width = Cm(34)
        desc_height = Cm(0.6)
        
        desc_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, desc_left, desc_top, desc_width, desc_height
        )
        desc_shape.fill.solid()
        desc_shape.fill.fore_color.rgb = RGBColor(139, 0, 0)  # 深红色背景
        desc_shape.line.fill.background()  # 无边框
        desc_shape.shadow.inherit = False  # 无阴影
        
        desc_frame = desc_shape.text_frame
        desc_frame.text = slide_data['description']
        desc_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        desc_frame.paragraphs[0].font.name = '微软雅黑'
        desc_frame.paragraphs[0].font.size = Pt(7)
        desc_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 白色文字
        
        # 3. 底层模块
        modules = slide_data['modules']
        left_modules = modules[:3]  # 左边3个
        right_modules = modules[3:5]  # 右边2个
        
        # 左边模块 - 固定高度4厘米
        left_module_height = Cm(4.2)  # 固定高度4厘米
        current_y = 3.0  # 起始垂直位置（厘米）
        
        for i, module in enumerate(left_modules):
            # 标题矩形
            title_left = Cm(1)
            title_top = Cm(current_y)
            title_width = Cm(6)
            title_height = Cm(0.8)
            
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, title_left, title_top, title_width, title_height
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = RGBColor(139, 0, 0)  # 深红色背景
            title_shape.line.fill.background()  # 无边框
            title_shape.shadow.inherit = False  # 无阴影
            
            title_frame = title_shape.text_frame
            title_frame.text = module['title']
            title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            title_frame.paragraphs[0].font.name = '微软雅黑'
            title_frame.paragraphs[0].font.size = Pt(14)
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 白色文字
            
            # 内容文本框 - 固定高度4厘米
            content_left = Cm(1)
            content_top = Cm(current_y + 0.8)
            content_width = Cm(14.5)
            content_height = left_module_height  # 固定高度4厘米
            
            content_shape = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_shape.shadow.inherit = False  # 无阴影
            content_frame = content_shape.text_frame
            content_frame.word_wrap = True
            
            # 应用格式化文本
            apply_text_formatting(content_frame, module['content'])
            
            # 设置字体
            for paragraph in content_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = '微软雅黑'
                    run.font.size = Pt(12)
            
            # 计算下一个模块的位置 - 使用固定间距
            current_y += 4.8 + 0.3  # 标题高度(0.8cm) + 内容高度(4cm) + 模块间距(0.3cm)
        
        # 右边模块 - 固定高度6厘米
        right_module_height = Cm(6.0)  # 固定高度6厘米
        current_y = 3.0  # 重置起始位置
        
        for i, module in enumerate(right_modules):
            # 标题矩形
            title_left = Cm(16)
            title_top = Cm(current_y)
            title_width = Cm(6)
            title_height = Cm(0.8)
            
            title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, title_left, title_top, title_width, title_height
            )
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = RGBColor(139, 0, 0)  # 深红色背景
            title_shape.line.fill.background()  # 无边框
            title_shape.shadow.inherit = False  # 无阴影
            
            title_frame = title_shape.text_frame
            title_frame.text = module['title']
            title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            title_frame.paragraphs[0].font.name = '微软雅黑'
            title_frame.paragraphs[0].font.size = Pt(14)
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 白色文字
            
            # 内容文本框 - 固定高度6厘米
            content_left = Cm(16)
            content_top = Cm(current_y + 0.8)
            content_width = Cm(14.5)
            content_height = right_module_height  # 固定高度6厘米
            
            content_shape = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_shape.shadow.inherit = False  # 无阴影
            content_frame = content_shape.text_frame
            content_frame.word_wrap = True
            
            # 应用格式化文本
            apply_text_formatting(content_frame, module['content'])
            
            # 设置字体
            for paragraph in content_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = '微软雅黑'
                    run.font.size = Pt(12)
            
            # 计算下一个模块的位置 - 使用固定间距
            current_y += 6.8 + 0.3  # 标题高度(0.8cm) + 内容高度(6cm) + 模块间距(0.3cm)
    
    # 保存文件
    prs.save(output_file)
    print(f"PPT已生成：{output_file}")


def markdown2pptx(markdown_file: str = "sample.md"):
    """主函数"""
    output_file = "output.pptx"
    
    try:
        create_pptx_from_markdown(markdown_file, output_file)
        print(f"已生成 {output_file}")
    except FileNotFoundError:
        print(f"错误：找不到文件 {markdown_file}")
        print("请确保在当前目录下存在 sample.md 文件")
    except Exception as e:
        print(f"转换过程中出现错误：{e}")


if __name__ == "__main__":
    markdown2pptx()